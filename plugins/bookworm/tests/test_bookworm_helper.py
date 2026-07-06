from __future__ import annotations

import subprocess
import sys
import tempfile
import unittest
from io import BytesIO
import importlib.util
import json
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from bookworm_helper import (
    citation_inventory,
    handoff_refined_note,
    refine_markdown,
    resolve_citation_markers,
    resolve_numeric_citations,
    source_counts,
    convert_refine_input,
    default_vault_roots,
    detect_vaults,
)


class RefineMarkdownTests(unittest.TestCase):
    def test_detects_icloud_obsidian_vault_once(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            home = Path(directory)
            brain = home / "Library" / "Mobile Documents" / "iCloud~md~obsidian" / "Documents" / "Brain"
            (brain / ".obsidian").mkdir(parents=True)
            (home / ".Trash" / "Old Vault" / ".obsidian").mkdir(parents=True)

            vaults = detect_vaults(default_vault_roots(home))

            self.assertEqual(vaults, [{"path": str(brain), "name": "Brain"}])

    def test_citation_inventory_keeps_marker_context_before_cleanup(self) -> None:
        first = "citeturn0search1"
        second = "citeturn1search2"
        source = (
            f"First claim about a mechanic {first}.\n\n"
            f"Second claim is separate {second}.\n"
        )

        self.assertEqual(
            citation_inventory(source),
            [
                {"line": 1, "marker": first, "context": f"First claim about a mechanic {first}."},
                {"line": 3, "marker": second, "context": f"Second claim is separate {second}."},
            ],
        )

    def test_preserves_ordinary_grouped_sources_without_numeric_citations(self) -> None:
        source = """## Тезисы

Текст исследования.

## Источники

### Официальные

- [Документация](https://example.com/docs)
- [Правила](https://example.com/rules)
"""
        result, report = resolve_numeric_citations(source)
        self.assertEqual(result, source)
        self.assertEqual(report["numeric_citations_scanned"], 0)

    def test_handoff_requires_confirmation_and_uses_title_filename(self) -> None:
        source = "# Принципы хороших интерфейсов\n\nИсходный текст.\n"

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source_path = root / "Inbox" / "deep-research-report.md"
            refined_path = root / "scratch" / "deep-research-report.md"
            destination_dir = root / "Library"
            source_path.parent.mkdir()
            refined_path.parent.mkdir()
            source_path.write_text(source, encoding="utf-8")
            refined_path.write_text("## Содержание\n", encoding="utf-8")
            (refined_path.parent / "manifest.json").write_text("{}", encoding="utf-8")

            with self.assertRaises(PermissionError):
                handoff_refined_note(
                    source_path,
                    refined_path,
                    destination_dir,
                    confirmation=None,
                    run_dir=refined_path.parent,
                )

            self.assertTrue(source_path.exists())
            self.assertTrue(refined_path.exists())

            destination = handoff_refined_note(
                source_path,
                refined_path,
                destination_dir,
                confirmation="user-confirmed",
                run_dir=refined_path.parent,
            )

            self.assertEqual(destination.name, "Принципы хороших интерфейсов.md")
            self.assertTrue(destination.exists())
            self.assertFalse(source_path.exists())
            self.assertFalse(refined_path.parent.exists())

    def test_handoff_copies_assets_to_shared_library_assets_folder(self) -> None:
        source = "# Принципы интерфейсов\n\nИсходный текст.\n"

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source_path = root / "Inbox" / "report.md"
            run_dir = root / "scratch"
            refined_path = run_dir / "refined.md"
            assets_dir = run_dir / "assets"
            destination_dir = root / "Vault" / "Library"
            source_path.parent.mkdir()
            assets_dir.mkdir(parents=True)
            source_path.write_text(source, encoding="utf-8")
            refined_path.write_text("## Содержание\n", encoding="utf-8")
            (assets_dir / "figure-001.png").write_bytes(b"png")

            destination = handoff_refined_note(
                source_path,
                refined_path,
                destination_dir,
                confirmation="user-confirmed",
                run_dir=run_dir,
                assets_dir=assets_dir,
            )

            self.assertTrue(destination.exists())
            self.assertTrue(
                (destination_dir / "assets" / "printsipy-interfeysov" / "figure-001.png").exists()
            )
            self.assertFalse(run_dir.exists())

    def test_handoff_ignores_empty_assets_directory(self) -> None:
        source = "# Empty assets\n\nText.\n"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source_path = root / "Inbox" / "source.md"
            run_dir = root / "scratch"
            refined_path = run_dir / "refined.md"
            assets_dir = run_dir / "assets"
            destination_dir = root / "Vault" / "Library"
            source_path.parent.mkdir()
            assets_dir.mkdir(parents=True)
            source_path.write_text(source, encoding="utf-8")
            refined_path.write_text("## Contents\n", encoding="utf-8")

            destination = handoff_refined_note(
                source_path,
                refined_path,
                destination_dir,
                confirmation="user-confirmed",
                run_dir=run_dir,
                assets_dir=assets_dir,
            )

            self.assertTrue(destination.exists())
            self.assertFalse((destination_dir / "assets" / "empty-assets").exists())
            self.assertFalse(run_dir.exists())

    def test_converts_markdown_to_temp_copy_without_changing_original(self) -> None:
        source = "# Title\n\nBody.\n"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            input_path = root / "source.md"
            output_path = root / "run" / "converted.md"
            input_path.write_text(source, encoding="utf-8")

            result = convert_refine_input(input_path, output_path)

            self.assertEqual(input_path.read_text(encoding="utf-8"), source)
            self.assertEqual(output_path.read_text(encoding="utf-8"), source)
            self.assertEqual(result["kind"], "markdown")
            self.assertEqual(result["assets"], [])

    def test_rejects_unsupported_refine_input_without_creating_note(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            input_path = root / "source.txt"
            output_path = root / "run" / "converted.md"
            input_path.write_text("plain text", encoding="utf-8")

            with self.assertRaisesRegex(ValueError, "Unsupported Refine input"):
                convert_refine_input(input_path, output_path)

            self.assertFalse(output_path.exists())

    @unittest.skipUnless(importlib.util.find_spec("docx"), "requires bundled python-docx")
    def test_converts_docx_structure_and_embeds_extracted_image(self) -> None:
        from docx import Document

        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\x0dIDAT\x08\xd7c\xf8\xcf\xc0\xf0\x1f\x00\x05\x00\x01\xff\x89\x99=\x1d\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "input.docx"
            output = root / "run" / "note.md"
            document = Document()
            document.add_heading("Document title", level=1)
            document.add_paragraph("A paragraph")
            document.add_paragraph("A list item", style="List Bullet")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Key"
            table.cell(0, 1).text = "Value"
            table.cell(1, 0).text = "A"
            table.cell(1, 1).text = "B"
            document.add_picture(BytesIO(png))
            document.save(source)

            result = convert_refine_input(source, output)

            rendered = output.read_text(encoding="utf-8")
            self.assertIn("# Document title", rendered)
            self.assertIn("- A list item", rendered)
            self.assertIn("| Key | Value |", rendered)
            self.assertEqual(len(result["assets"]), 1)
            self.assertTrue(Path(result["assets"][0]).exists())

    @unittest.skipUnless(importlib.util.find_spec("docx"), "requires bundled python-docx")
    def test_preserves_docx_hyperlinks_as_markdown_title_links(self) -> None:
        from docx import Document
        from docx.enum.style import WD_STYLE_TYPE
        from docx.opc.constants import RELATIONSHIP_TYPE
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "input.docx"
            output = root / "run" / "note.md"
            document = Document()
            paragraph = document.add_paragraph("Read ")
            relationship_id = paragraph.part.relate_to(
                "https://example.com/primary-source",
                RELATIONSHIP_TYPE.HYPERLINK,
                is_external=True,
            )
            hyperlink = OxmlElement("w:hyperlink")
            hyperlink.set(qn("r:id"), relationship_id)
            run = OxmlElement("w:r")
            text = OxmlElement("w:t")
            text.text = "Primary source"
            run.append(text)
            hyperlink.append(run)
            paragraph._p.append(hyperlink)
            document.save(source)

            convert_refine_input(source, output)

            self.assertIn(
                "[Primary source](https://example.com/primary-source)",
                output.read_text(encoding="utf-8"),
            )

    @unittest.skipUnless(importlib.util.find_spec("pptx"), "requires bundled python-pptx")
    def test_converts_pptx_text_notes_and_illustrations(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\x0dIDAT\x08\xd7c\xf8\xcf\xc0\xf0\x1f\x00\x05\x00\x01\xff\x89\x99=\x1d\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "slides.pptx"
            output = root / "run" / "note.md"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = "Slide claim"
            slide.shapes.add_picture(BytesIO(png), Inches(1), Inches(2))
            presentation.save(source)

            result = convert_refine_input(source, output)

            self.assertIn("## Slide 1", output.read_text(encoding="utf-8"))
            self.assertIn("Slide claim", output.read_text(encoding="utf-8"))
            self.assertEqual(len(result["assets"]), 1)

    @unittest.skipUnless(importlib.util.find_spec("pdfplumber"), "requires bundled PDF libraries")
    def test_converts_pdf_text_and_embedded_image(self) -> None:
        from reportlab.pdfgen.canvas import Canvas

        png = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\x0dIDAT\x08\xd7c\xf8\xcf\xc0\xf0\x1f\x00\x05\x00\x01\xff\x89\x99=\x1d\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pdf"
            image = root / "image.png"
            output = root / "run" / "note.md"
            image.write_bytes(png)
            canvas = Canvas(str(source))
            canvas.drawString(72, 720, "PDF claim")
            canvas.drawImage(str(image), 72, 600, width=20, height=20)
            canvas.save()

            result = convert_refine_input(source, output)

            self.assertIn("PDF claim", output.read_text(encoding="utf-8"))
            self.assertEqual(len(result["assets"]), 1)

    def test_removes_raw_citation_markers_from_reader_output(self) -> None:
        citation = "citeturn0search1"
        source = f"# Title\n\nClaim supported only by {citation}.\n\n## Findings\n"

        result = refine_markdown(source, toc_title="Содержание")

        self.assertNotIn(citation, result)

    def test_removes_citations_preserves_urls_and_adds_toc(self) -> None:
        source = """# Board Game Mechanics

Intro citeturn0search1 with a [source](https://example.com/research).

## Findings

Useful text.

## Sources

- https://example.com/research
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertNotIn("cite", result)
        self.assertNotIn("# Board Game Mechanics", result)
        self.assertIn("[source](https://example.com/research)", result)
        self.assertIn("## Содержание", result)
        self.assertIn("- [[#Findings|Findings]]", result)
        self.assertIn("- [[#Sources|Sources]]", result)

    def test_replaces_legacy_markdown_anchor_toc_without_duplicates(self) -> None:
        source = """# Title

## Содержание

- [Findings](#findings)

## Findings

Useful text.
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertEqual(result.count("## Содержание"), 1)
        self.assertNotIn("- [Findings](#findings)", result)
        self.assertIn("- [[#Findings|Findings]]", result)

    def test_replaces_legacy_obsidian_toc_without_duplicates(self) -> None:
        source = """# Заголовок

## Содержание

- [[#Находки|Находки]]

## Находки

Полезный текст.
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertEqual(result.count("## Содержание"), 1)
        self.assertEqual(result.count("- [[#Находки|Находки]]"), 1)

    def test_turns_two_related_label_value_fields_into_a_table(self) -> None:
        source = """## Профиль

**Определение:** Игрок размещает рабочих на общих ячейках.

**Сильные стороны:** Блокирование и читаемое планирование.

**Игры:** [Agricola](https://example.com/agricola)
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn("| Параметр | Описание |", result)
        self.assertIn("| Определение | Игрок размещает рабочих на общих ячейках. |", result)
        self.assertIn("| Игры | [Agricola](https://example.com/agricola) |", result)
        self.assertNotIn("**Сильные стороны:**", result)

    def test_turns_bold_label_with_external_colon_into_a_table(self) -> None:
        source = """## Профиль

**Автор**: Francesco Cirillo.
**Базовая единица**: Фокус-сессия и отношение ко времени.
**Ограничения**: Не является полной системой проектов.
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn("| Параметр | Описание |", result)
        self.assertIn("| Автор | Francesco Cirillo. |", result)
        self.assertIn("| Ограничения | Не является полной системой проектов. |", result)

    def test_keeps_prose_heavy_label_value_runs_out_of_tables(self) -> None:
        source = """## Профиль

**Кратко:** Это длинное описание механики, которое объясняет её устройство, варианты применения и последствия для проектирования. Оно намеренно содержит больше одного предложения, чтобы оставаться читабельным как обычный текст, а не как узкая ячейка таблицы.

**Проектирование:** Второй развёрнутый абзац описывает ограничения, типовые ошибки и способы баланса. Он также достаточно длинный, чтобы таблица ухудшила сканирование заметки.
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertNotIn("| Параметр | Описание |", result)
        self.assertIn("**Кратко:**", result)
        self.assertIn("**Проектирование:**", result)

    def test_localizes_existing_parameter_table_headers(self) -> None:
        source = """## Профиль

| Parameter | Description |
| --- | --- |
| Размер руки | Число карт |
"""
        result = refine_markdown(source, toc_title="Содержание")
        self.assertIn("| Параметр | Описание |", result)
        self.assertNotIn("| Parameter | Description |", result)

    def test_localizes_generic_english_heading_in_russian_research(self) -> None:
        result = refine_markdown("## Executive summary\n\nРусский текст.\n")
        self.assertIn("## Исполнительное резюме", result)

    def test_localizes_common_english_research_structure_headings_in_russian_note(self) -> None:
        source = """## Page structure and story flow

Русский текст про лендинг и wireframe.

## Layout grids, spacing and typography

Ещё русский текст.
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn("[[#Структура страницы и сценарий|Структура страницы и сценарий]]", result)
        self.assertIn("[[#Сетки, интервалы и типографика|Сетки, интервалы и типографика]]", result)
        self.assertIn("## Структура страницы и сценарий", result)
        self.assertIn("## Сетки, интервалы и типографика", result)
        self.assertNotIn("Page structure and story flow", result)
        self.assertNotIn("Layout grids, spacing and typography", result)

    def test_preserves_established_terms_inside_localized_headings(self) -> None:
        source = """## AI-agent instruction block

Русский текст для Codex.

## Ready for visual design checklist

Русский текст про visual design.
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn("## Инструкция для AI-агента", result)
        self.assertIn("## Чеклист готовности к visual design", result)
        self.assertIn("[[#Инструкция для AI-агента|Инструкция для AI-агента]]", result)

    def test_escapes_pipes_inside_existing_markdown_table_links(self) -> None:
        source = """## Сравнение

| Метод | Источник |
| --- | --- |
| GTD | [What is GTD? | Getting Things Done](https://gettingthingsdone.com/) |
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn(
            "[What is GTD? \\| Getting Things Done](https://gettingthingsdone.com/)",
            result,
        )
        self.assertNotIn(
            "[What is GTD? | Getting Things Done](https://gettingthingsdone.com/)",
            result,
        )

    def test_compacts_mermaid_without_rasterizing_it(self) -> None:
        source = """## Процесс

```mermaid
flowchart TB
    A[Начало] --> B[Следующий шаг]
```
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn("```mermaid", result)
        self.assertIn('%%{init: {"flowchart": {"useMaxWidth": false, "nodeSpacing": 20, "rankSpacing": 25}} }%%', result)
        self.assertNotIn("![[assets/", result)

    def test_verticalizes_horizontal_mermaid_without_splitting_the_graph(self) -> None:
        source = """## Поток

```mermaid
flowchart LR
    A[Входящее] --> B[Разобрать] --> C[Следующее действие]
```
"""

        result = refine_markdown(source, toc_title="Содержание")

        self.assertIn("flowchart TD", result)
        self.assertNotIn("flowchart LR", result)
        self.assertIn("A[Входящее] --> B[Разобрать] --> C[Следующее действие]", result)
        self.assertEqual(result.count("```mermaid"), 1)

    def test_counts_each_source_bearing_construct(self) -> None:
        source = """Claim citeturn0search1 [named source](https://example.com/a).

Bare source: https://example.org/b.

Footnote reference[^source].
"""

        self.assertEqual(
            source_counts(source),
            {
                "citation_markers": 1,
                "markdown_links": 1,
                "bare_urls": 1,
                "footnote_references": 1,
            },
        )

    def test_inserts_verified_title_link_at_raw_citation_claim(self) -> None:
        marker = "citeturn0search1"
        source = f"A mechanic is documented by BoardGameGeek {marker}."

        result, report = resolve_citation_markers(
            source,
            {
                marker: {
                    "title": "BoardGameGeek",
                    "url": "https://boardgamegeek.com/",
                }
            },
        )

        self.assertEqual(
            result,
            "A mechanic is documented by BoardGameGeek [BoardGameGeek](https://boardgamegeek.com/).",
        )
        self.assertEqual(
            report,
            {"markers_scanned": 1, "verified_title_links_inserted": 1, "unresolved": 0},
        )

    def test_removes_unverified_marker_and_reports_it_without_inventing_link(self) -> None:
        marker = "citeturn0search1"
        result, report = resolve_citation_markers("Claim " + marker + ".", {})

        self.assertEqual(result, "Claim .")
        self.assertEqual(
            report,
            {"markers_scanned": 1, "verified_title_links_inserted": 0, "unresolved": 1},
        )

    def test_replaces_numeric_citations_with_readable_sources_section(self) -> None:
        source = """## Mechanics

Worker placement creates scarcity [1]. It can pair with deckbuilding [2].

## Sources

[1] [Worker Placement | BoardGameGeek](https://boardgamegeek.com/browse/boardgamemechanic)
[2] [Dune: Imperium](https://duneimperium.com/)
"""

        result, report = resolve_numeric_citations(source)

        self.assertNotIn("[1]", result)
        self.assertNotIn("[2]", result)
        self.assertIn("## Sources", result)
        self.assertIn("- [Worker Placement | BoardGameGeek](https://boardgamegeek.com/browse/boardgamemechanic)", result)
        self.assertIn("- [Dune: Imperium](https://duneimperium.com/)", result)
        self.assertEqual(report, {"numeric_citations_scanned": 2, "numeric_sources_resolved": 2, "numeric_unresolved": 0})

    def test_uses_verified_mapping_for_numbered_source_without_url(self) -> None:
        source = """Claim about a mechanism [66].

## Источники

[66] Worker Placement | Board Game Mechanic
"""

        result, report = resolve_numeric_citations(
            source,
            {"66": {"title": "Worker Placement | BoardGameGeek", "url": "https://boardgamegeek.com/browse/boardgamemechanic"}},
        )

        self.assertNotIn("[66]", result)
        self.assertIn("- [Worker Placement | BoardGameGeek](https://boardgamegeek.com/browse/boardgamemechanic)", result)
        self.assertEqual(report["numeric_unresolved"], 0)

    def test_removes_unmapped_numeric_citation_and_reports_unresolved(self) -> None:
        source = """Claim [9].

## Sources

[9] Unknown article
"""

        result, report = resolve_numeric_citations(source)

        self.assertNotIn("[9]", result)
        self.assertNotIn("Unknown article", result)
        self.assertEqual(report, {"numeric_citations_scanned": 1, "numeric_sources_resolved": 0, "numeric_unresolved": 1})

    def test_cli_writes_refined_copy_without_changing_source(self) -> None:
        source = "# Title\n\nText citeturn1search2\n\n## Section\n"

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            input_path = root / "research.md"
            output_path = root / "refined.md"
            input_path.write_text(source, encoding="utf-8")

            completed = subprocess.run(
                [
                    sys.executable,
                    str(ROOT / "scripts" / "bookworm_helper.py"),
                    "refine-markdown",
                    str(input_path),
                    "--out",
                    str(output_path),
                ],
                check=True,
                capture_output=True,
                text=True,
            )

            self.assertEqual(input_path.read_text(encoding="utf-8"), source)
            self.assertTrue(output_path.exists())
            self.assertNotIn("cite", output_path.read_text(encoding="utf-8"))
            self.assertIn(str(output_path), completed.stdout)
            self.assertIn('"suggested_filename": "Title.md"', completed.stdout)

    def test_cli_reports_verified_and_unresolved_citation_markers(self) -> None:
        marker = "citeturn1search2"
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            input_path = root / "research.md"
            output_path = root / "refined.md"
            verified_path = root / "verified.json"
            input_path.write_text(f"# Title\n\nClaim {marker}.\n", encoding="utf-8")
            verified_path.write_text(
                json.dumps({marker: {"title": "Source", "url": "https://example.com/"}}),
                encoding="utf-8",
            )

            completed = subprocess.run(
                [
                    sys.executable,
                    str(ROOT / "scripts" / "bookworm_helper.py"),
                    "refine-markdown",
                    str(input_path),
                    "--out",
                    str(output_path),
                    "--verified-sources",
                    str(verified_path),
                ],
                check=True,
                capture_output=True,
                text=True,
            )

            report = json.loads(completed.stdout)
            self.assertEqual(report["markers_scanned"], 1)
            self.assertEqual(report["verified_title_links_inserted"], 1)
            self.assertEqual(report["unresolved"], 0)
            self.assertIn("[Source](https://example.com/)", output_path.read_text(encoding="utf-8"))

    def test_cli_normalizes_numbered_sources_and_reports_numeric_counts(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            input_path = root / "research.md"
            output_path = root / "refined.md"
            input_path.write_text(
                "## Claim\n\nFact [1].\n\n## Sources\n\n[1] [Primary source](https://example.com/)\n",
                encoding="utf-8",
            )

            completed = subprocess.run(
                [
                    sys.executable,
                    str(ROOT / "scripts" / "bookworm_helper.py"),
                    "refine-markdown",
                    str(input_path),
                    "--out",
                    str(output_path),
                ],
                check=True,
                capture_output=True,
                text=True,
            )

            report = json.loads(completed.stdout)
            refined = output_path.read_text(encoding="utf-8")
            self.assertNotIn("[1]", refined)
            self.assertIn("- [Primary source](https://example.com/)", refined)
            self.assertEqual(report["numeric_citations_scanned"], 1)
            self.assertEqual(report["numeric_sources_resolved"], 1)
            self.assertEqual(report["numeric_unresolved"], 0)


if __name__ == "__main__":
    unittest.main()
